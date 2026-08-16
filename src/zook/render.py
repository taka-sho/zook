"""PPTX rendering.

Decisions carried over from docs/detailed-design-pptx.md (confirmed via the
prototype in prototype/build_prototype.py):

- 1 logical unit == 9525 EMU == 1px @ 96dpi, for both supported aspect ratios
  (sec8.5).
- Groups get absolute (slide-level) coordinates on every descendant; relying
  on python-pptx's automatic chOff/chExt = off/ext extent recalculation
  (sec8.4) instead of a hand-rolled helper.
- Connector connection-point indices: 0=top, 1=left, 2=bottom, 3=right
  (sec8.2), chosen from the relative position of the two endpoints.
- Connector labels cannot be injected as cxnSp/txBody (invalid per the OOXML
  schema); they are separate textboxes at the connector midpoint (sec8.3).
- Container labels need an explicit TOP vertical anchor or python-pptx's
  default vertical-centering makes them collide with nested content.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .layout import (
    Box,
    content_offset,
    is_shape_node,
    label_box_height,
    label_gap,
    link_label_anchor,
    link_label_size,
    link_render_plan,
    node_label_font_size,
    resolve_container_style,
)
from .model import Diagram, Link
from .registry import MultiRegistry

LOGICAL_TO_EMU = 9525
LOGICAL_TO_PT = LOGICAL_TO_EMU / 12700  # 1pt = 12700 EMU
CORNER_BADGE_SIZE = 20  # logical units; corner icon for a container's group style (e.g. "AWS Cloud")
CORNER_BADGE_PADDING = 6


def E(value: float) -> Emu:
    return Emu(int(round(value * LOGICAL_TO_EMU)))


def _set_dashed(line) -> None:
    ln = line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))


def _apply_label_position(text_frame, position: str) -> None:
    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM if "bottom" in position else MSO_ANCHOR.TOP
    align = PP_ALIGN.CENTER if "center" in position else PP_ALIGN.LEFT
    for p in text_frame.paragraphs:
        p.alignment = align


def _add_container_rect(shapes, box: Box, registry: MultiRegistry):
    element = box.element
    style = resolve_container_style(element, registry)

    rect = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(box.abs_x), E(box.abs_y), E(box.width), E(box.height))
    if style.fill_color:
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(style.fill_color.lstrip("#"))
    else:
        rect.fill.background()
    rect.line.color.rgb = RGBColor.from_string(style.border_color.lstrip("#"))
    rect.line.width = Pt(style.border_width)
    if style.dashed:
        _set_dashed(rect.line)
    rect.shadow.inherit = False

    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Pt(4)
    tf.text = style.label_text
    _apply_label_position(tf, style.label_position)
    for p in tf.paragraphs:
        p.font.size = Pt(style.label_font_size)
        p.font.color.rgb = RGBColor.from_string(style.border_color.lstrip("#"))

    # Corner badge (e.g. the AWS Cloud logo) marks visually where a boundary
    # like "AWS Cloud" starts, per groups.<type>.icon in the icon registry.
    if style.corner_icon and style.corner_icon.exists() and "left" in style.label_position:
        badge_x = box.abs_x + CORNER_BADGE_PADDING
        badge_y = (
            box.abs_y + CORNER_BADGE_PADDING
            if "top" in style.label_position
            else box.abs_y + box.height - CORNER_BADGE_PADDING - CORNER_BADGE_SIZE
        )
        shapes.add_picture(
            str(style.corner_icon), E(badge_x), E(badge_y), E(CORNER_BADGE_SIZE), E(CORNER_BADGE_SIZE)
        )
        tf.margin_left = Pt((CORNER_BADGE_SIZE + CORNER_BADGE_PADDING) * LOGICAL_TO_PT)

    return rect


def _add_node_label(shapes, box: Box, text: str, position: str) -> None:
    gap = label_gap(box.element)
    font_size = node_label_font_size(box.element)
    box_height = label_box_height(font_size)
    dx, dy = content_offset(box)
    footprint_x, footprint_y = box.abs_x - dx, box.abs_y - dy
    if position == "below":
        tb = shapes.add_textbox(E(footprint_x), E(box.abs_y + box.height + gap), E(box.footprint_w), E(box_height))
        align = PP_ALIGN.CENTER
    elif position == "above":
        tb = shapes.add_textbox(E(footprint_x), E(footprint_y), E(box.footprint_w), E(box_height))
        align = PP_ALIGN.CENTER
    else:  # right
        tb = shapes.add_textbox(
            E(box.abs_x + box.width + gap),
            E(box.abs_y),
            E(max(box.footprint_w - box.width - gap, 60)),
            E(box.height),
        )
        align = PP_ALIGN.LEFT
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = align
    tf.paragraphs[0].font.size = Pt(font_size)


_SHAPE_MSO = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "circle": MSO_SHAPE.OVAL,
}


def _add_shape_node(shapes, box: Box):
    element = box.element
    shape = shapes.add_shape(_SHAPE_MSO[element.style["shape"]], E(box.abs_x), E(box.abs_y), E(box.width), E(box.height))

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(element.style.get("fillColor", "FFFFFF").lstrip("#").upper())
    shape.line.color.rgb = RGBColor.from_string(element.style.get("borderColor", "000000").lstrip("#").upper())

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    label_text = element.label if element.label is not None else element.type
    tf.paragraphs[0].text = label_text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(node_label_font_size(element))
    tf.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    return shape


def _add_node(shapes, box: Box, registry: MultiRegistry):
    # Unresolved type / missing icon file are reported by
    # layout.icon_resolution_warnings() (run once, before rendering, so
    # `validate` sees the same warnings `build` would) - fall back to the
    # placeholder silently here rather than reporting it a second time.
    element = box.element
    if is_shape_node(element):
        return _add_shape_node(shapes, box)

    icon_entry = registry.resolve_icon(element.type, element.provider)
    if icon_entry is None or not icon_entry.file.exists():
        icon_path = registry.placeholder_icon
    else:
        icon_path = icon_entry.file

    pic = shapes.add_picture(str(icon_path), E(box.abs_x), E(box.abs_y), E(box.width), E(box.height))

    label_position = element.style.get("labelPosition", "below")
    if label_position != "none":
        label_text = element.label if element.label is not None else (
            icon_entry.label if icon_entry and icon_entry.label else element.type
        )
        _add_node_label(shapes, box, label_text, label_position)
    return pic


def _render_element(shapes, box: Box, registry: MultiRegistry, shape_index: dict) -> None:
    element = box.element
    if element.is_container:
        group = shapes.add_group_shape()
        group_shapes = group.shapes
        rect = _add_container_rect(group_shapes, box, registry)
        shape_index[element.id] = (rect, box)
        for child in box.children:
            _render_element(group_shapes, child, registry, shape_index)
    else:
        pic = _add_node(shapes, box, registry)
        shape_index[element.id] = (pic, box)


def _add_link_label(shapes, midpoint: tuple[float, float], text: str, font_size: float) -> None:
    """sec8.3: p:cxnSp cannot carry txBody; use a separate midpoint textbox.

    `midpoint` is in logical units - the chord midpoint of the path's two
    endpoints, matching exactly where link_crossing_warnings() predicts the
    label sits, so render and the overlap check never disagree."""
    mx, my = E(midpoint[0]), E(midpoint[1])
    label_w, label_h = link_label_size(font_size)
    w, h = E(label_w), E(label_h)
    tb = shapes.add_textbox(Emu(int(mx - w / 2)), Emu(int(my - h / 2)), w, h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(font_size)


_CONNECTOR_TYPES = {"straight": MSO_CONNECTOR.STRAIGHT, "elbow": MSO_CONNECTOR.ELBOW, "curved": MSO_CONNECTOR.CURVE}


_NON_RECT_SHAPES = {"diamond", "circle"}


def _is_non_rect_shape_node(box: Box) -> bool:
    # diamond/circle prstGeom connection sites aren't laid out top/left/
    # bottom/right in bounding-box order the way rect/rounded/pictures/
    # containers are - PowerPoint/LibreOffice snap a glued (stCxn/endCxn)
    # connector to the shape's own site for that index, silently overriding
    # the literal begin_x/y we write below (python-pptx's begin_connect/
    # end_connect docstring warns of exactly this for non-rectangular
    # shapes). Skipping the glue for these keeps the literal, geometrically
    # correct point instead.
    return is_shape_node(box.element) and box.element.style.get("shape") in _NON_RECT_SHAPES


def _style_connector(conn) -> None:
    conn.line.color.rgb = RGBColor.from_string("545B64")
    conn.line.width = Pt(1.25)


def _add_arrowhead(conn, tag: str) -> None:
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn(tag), {"type": "triangle", "w": "med", "len": "med"}))


def _render_polyline_link(shapes, link: Link, path, from_shape, to_shape, start_idx, end_idx, from_box, to_box):
    """A waypoint link renders as one straight connector per segment (spike
    finding: the lowest-risk way to draw an N-bend polyline - it reuses the
    exact begin/end-override primitive the single-segment case uses). The
    arrowhead goes on the last segment only (and the head, for `both`, on the
    first) so the chain reads as one arrow."""
    segments = []
    for (x0, y0), (x1, y1) in zip(path, path[1:]):
        conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(x0), E(y0), E(x1), E(y1))
        conn.begin_x, conn.begin_y = E(x0), E(y0)
        conn.end_x, conn.end_y = E(x1), E(y1)
        _style_connector(conn)
        segments.append(conn)

    # Glue the two ends to their shapes (interior joints stay at the explicit
    # via points); re-assert the exact endpoints, since begin/end_connect snap.
    if not _is_non_rect_shape_node(from_box):
        segments[0].begin_connect(from_shape, start_idx)
        segments[0].begin_x, segments[0].begin_y = E(path[0][0]), E(path[0][1])
    if not _is_non_rect_shape_node(to_box):
        segments[-1].end_connect(to_shape, end_idx)
        segments[-1].end_x, segments[-1].end_y = E(path[-1][0]), E(path[-1][1])

    if link.arrow != "none":
        _add_arrowhead(segments[-1], "a:tailEnd")
        if link.arrow == "both":
            _add_arrowhead(segments[0], "a:headEnd")


def _render_link(shapes, link: Link, shape_index: dict) -> None:
    from_shape, from_box = shape_index[link.from_id]
    to_shape, to_box = shape_index[link.to_id]
    start_idx, end_idx, eff_style, path = link_render_plan(from_box, to_box, link)

    if link.waypoints:
        _render_polyline_link(shapes, link, path, from_shape, to_shape, start_idx, end_idx, from_box, to_box)
    else:
        conn = shapes.add_connector(_CONNECTOR_TYPES[eff_style], E(0), E(0), E(1), E(1))
        if not _is_non_rect_shape_node(from_box):
            conn.begin_connect(from_shape, start_idx)
        if not _is_non_rect_shape_node(to_box):
            conn.end_connect(to_shape, end_idx)
        # begin_connect/end_connect snap to the connected shape's own edge; override
        # with our (possibly label-aware, sec "connection_point") points so a
        # bottom/top exit past a below/above label, and the auto-elbow bend for a
        # diagonal `straight` link, both actually render as planned.
        conn.begin_x, conn.begin_y = E(path[0][0]), E(path[0][1])
        conn.end_x, conn.end_y = E(path[-1][0]), E(path[-1][1])
        _style_connector(conn)
        if link.arrow != "none":
            _add_arrowhead(conn, "a:tailEnd")
            if link.arrow == "both":
                _add_arrowhead(conn, "a:headEnd")

    if link.label:
        _add_link_label(shapes, link_label_anchor(path), link.label, link.label_font_size)


def render(diagram: Diagram, root_box: Box, registry: MultiRegistry) -> Presentation:
    prs = Presentation()
    canvas_w, canvas_h = diagram.canvas.size
    prs.slide_width = E(canvas_w)
    prs.slide_height = E(canvas_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if diagram.canvas.background:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(diagram.canvas.background.lstrip("#"))

    shape_index: dict[str, tuple] = {}
    for child_box in root_box.children:
        _render_element(slide.shapes, child_box, registry, shape_index)

    for link in diagram.links:
        _render_link(slide.shapes, link, shape_index)

    return prs
