"""PPTX rendering.

Decisions carried over from docs/detailed-design-pptx.md (confirmed via the
prototype in prototype/build_prototype.py):

- 1 logical unit == 9525 EMU == 1px @ 96dpi, for both supported aspect ratios
  (sec8.5).
- Groups get absolute (slide-level) coordinates on every descendant; relying
  on python-pptx's automatic chOff/chExt = off/ext extent recalculation
  (sec8.4) instead of a hand-rolled helper.
- Connector connection-point indices: 0=top, 1=left, 2=bottom, 3=right
  (sec8.2), chosen from the relative position of the two endpoints.
- Connector labels cannot be injected as cxnSp/txBody (invalid per the OOXML
  schema); they are separate textboxes at the connector midpoint (sec8.3).
- Container labels need an explicit TOP vertical anchor or python-pptx's
  default vertical-centering makes them collide with nested content.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .errors import Warnings
from .layout import LABEL_RESERVE, Box, content_offset
from .model import Diagram, Link
from .registry import Registry

LOGICAL_TO_EMU = 9525
LOGICAL_TO_PT = LOGICAL_TO_EMU / 12700  # 1pt = 12700 EMU
CORNER_BADGE_SIZE = 20  # logical units; corner icon for a container's group style (e.g. "AWS Cloud")
CORNER_BADGE_PADDING = 6


def E(value: float) -> Emu:
    return Emu(int(round(value * LOGICAL_TO_EMU)))


def _set_dashed(line) -> None:
    ln = line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))


def _apply_label_position(text_frame, position: str) -> None:
    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM if "bottom" in position else MSO_ANCHOR.TOP
    align = PP_ALIGN.CENTER if "center" in position else PP_ALIGN.LEFT
    for p in text_frame.paragraphs:
        p.alignment = align


def _add_container_rect(shapes, box: Box, registry: Registry):
    element = box.element
    group_style = registry.resolve_group(element.type)
    style = element.style or {}

    border_color = style.get("borderColor") or (group_style.border_color if group_style else "#5A6B86")
    fill_color = style.get("fillColor") or (group_style.fill_color if group_style else None)
    border_width = style.get("borderWidth", group_style.border_width if group_style else 1)
    dashed = group_style.dashed if group_style else False
    label_position = style.get("labelPosition") or (group_style.label_position if group_style else "top-left")
    label_text = element.label if element.label is not None else (group_style.label if group_style else "")

    rect = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(box.abs_x), E(box.abs_y), E(box.width), E(box.height))
    if fill_color:
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))
    else:
        rect.fill.background()
    rect.line.color.rgb = RGBColor.from_string(border_color.lstrip("#"))
    rect.line.width = Pt(border_width)
    if dashed:
        _set_dashed(rect.line)
    rect.shadow.inherit = False

    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Pt(4)
    tf.text = label_text
    _apply_label_position(tf, label_position)
    for p in tf.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor.from_string(border_color.lstrip("#"))

    # Corner badge (e.g. the AWS Cloud logo) marks visually where a boundary
    # like "AWS Cloud" starts, per groups.<type>.icon in the icon registry.
    if group_style and group_style.icon and group_style.icon.exists() and "left" in label_position:
        badge_x = box.abs_x + CORNER_BADGE_PADDING
        badge_y = (
            box.abs_y + CORNER_BADGE_PADDING
            if "top" in label_position
            else box.abs_y + box.height - CORNER_BADGE_PADDING - CORNER_BADGE_SIZE
        )
        shapes.add_picture(
            str(group_style.icon), E(badge_x), E(badge_y), E(CORNER_BADGE_SIZE), E(CORNER_BADGE_SIZE)
        )
        tf.margin_left = Pt((CORNER_BADGE_SIZE + CORNER_BADGE_PADDING) * LOGICAL_TO_PT)

    return rect


def _add_node_label(shapes, box: Box, text: str, position: str) -> None:
    dx, dy = content_offset(box)
    footprint_x, footprint_y = box.abs_x - dx, box.abs_y - dy
    if position == "below":
        tb = shapes.add_textbox(
            E(footprint_x), E(box.abs_y + box.height + 4), E(box.footprint_w), E(LABEL_RESERVE - 4)
        )
        align = PP_ALIGN.CENTER
    elif position == "above":
        tb = shapes.add_textbox(E(footprint_x), E(footprint_y), E(box.footprint_w), E(LABEL_RESERVE - 4))
        align = PP_ALIGN.CENTER
    else:  # right
        tb = shapes.add_textbox(
            E(box.abs_x + box.width + 6), E(box.abs_y), E(max(box.footprint_w - box.width - 6, 60)), E(box.height)
        )
        align = PP_ALIGN.LEFT
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = align
    tf.paragraphs[0].font.size = Pt(9)


def _add_node(shapes, box: Box, registry: Registry, warnings: Warnings):
    element = box.element
    icon_entry = registry.resolve_icon(element.type)
    if icon_entry is None:
        warnings.add(f"unknown type {element.type!r} for node {element.id!r}; using placeholder icon")
        icon_path = registry.placeholder_icon
    else:
        icon_path = icon_entry.file
        if not icon_path.exists():
            warnings.add(f"icon file missing for type {element.type!r} ({icon_path}); using placeholder icon")
            icon_path = registry.placeholder_icon

    pic = shapes.add_picture(str(icon_path), E(box.abs_x), E(box.abs_y), E(box.width), E(box.height))

    label_position = element.style.get("labelPosition", "below")
    if label_position != "none":
        label_text = element.label if element.label is not None else (
            icon_entry.label if icon_entry and icon_entry.label else element.type
        )
        _add_node_label(shapes, box, label_text, label_position)
    return pic


def _render_element(shapes, box: Box, registry: Registry, warnings: Warnings, shape_index: dict) -> None:
    element = box.element
    if element.is_container:
        group = shapes.add_group_shape()
        group_shapes = group.shapes
        rect = _add_container_rect(group_shapes, box, registry)
        shape_index[element.id] = (rect, box)
        for child in box.children:
            _render_element(group_shapes, child, registry, warnings, shape_index)
    else:
        pic = _add_node(shapes, box, registry, warnings)
        shape_index[element.id] = (pic, box)


def _choose_connection_indices(from_box: Box, to_box: Box) -> tuple[int, int]:
    """sec8.2: idx 0=top, 1=left, 2=bottom, 3=right. Pick the edge facing the other shape."""
    fcx, fcy = from_box.abs_x + from_box.width / 2, from_box.abs_y + from_box.height / 2
    tcx, tcy = to_box.abs_x + to_box.width / 2, to_box.abs_y + to_box.height / 2
    dx, dy = tcx - fcx, tcy - fcy
    if abs(dx) >= abs(dy):
        return (3, 1) if dx >= 0 else (1, 3)
    return (2, 0) if dy >= 0 else (0, 2)


def _add_link_label(shapes, conn, text: str) -> None:
    """sec8.3: p:cxnSp cannot carry txBody; use a separate midpoint textbox."""
    mx, my = (conn.begin_x + conn.end_x) / 2, (conn.begin_y + conn.end_y) / 2
    w, h = E(60), E(18)
    tb = shapes.add_textbox(Emu(int(mx - w / 2)), Emu(int(my - h / 2)), w, h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(8)


_CONNECTOR_TYPES = {"straight": MSO_CONNECTOR.STRAIGHT, "elbow": MSO_CONNECTOR.ELBOW, "curved": MSO_CONNECTOR.CURVE}


def _render_link(shapes, link: Link, shape_index: dict) -> None:
    from_shape, from_box = shape_index[link.from_id]
    to_shape, to_box = shape_index[link.to_id]
    start_idx, end_idx = _choose_connection_indices(from_box, to_box)

    conn = shapes.add_connector(_CONNECTOR_TYPES[link.style], E(0), E(0), E(1), E(1))
    conn.begin_connect(from_shape, start_idx)
    conn.end_connect(to_shape, end_idx)
    conn.line.color.rgb = RGBColor.from_string("545B64")
    conn.line.width = Pt(1.25)

    if link.arrow != "none":
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
        if link.arrow == "both":
            ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))

    if link.label:
        _add_link_label(shapes, conn, link.label)


def render(diagram: Diagram, root_box: Box, registry: Registry, warnings: Warnings) -> Presentation:
    prs = Presentation()
    canvas_w, canvas_h = diagram.canvas.size
    prs.slide_width = E(canvas_w)
    prs.slide_height = E(canvas_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if diagram.canvas.background:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(diagram.canvas.background.lstrip("#"))

    shape_index: dict[str, tuple] = {}
    for child_box in root_box.children:
        _render_element(slide.shapes, child_box, registry, warnings, shape_index)

    for link in diagram.links:
        _render_link(slide.shapes, link, shape_index)

    return prs
