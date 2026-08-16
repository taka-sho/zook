"""Tests for link waypoints: explicit polyline routing that detours a
connector around an obstacle it would otherwise cross.

Grounding invariant: the path is derived by link_render_plan(), and every
checker/renderer walks that path, so a waypoint that routes clear of an
element removes its crossing warning and the diagram builds/previews/exports
without special-casing.
"""

import pytest
from pptx.shapes.connector import Connector

from zook.drawio import export_drawio
from zook.errors import DiagramError
from zook.layout import (
    build_layout,
    link_crossing_warnings,
    link_label_anchor,
    link_render_plan,
)
from zook.model import parse_diagram
from zook.registry import load_registries
from zook.render import render
from zook.validate import validate

REGISTRY = load_registries()

# A over B over X on one vertical line: a plain A->X link crosses B.
def _doc(link_extra):
    link = {"from": "A", "to": "X"}
    link.update(link_extra)
    return {
        "version": "1.0",
        "canvas": {"aspectRatio": "16:9"},
        "elements": [
            {"kind": "node", "id": "A", "type": "EC2", "x": 300, "y": 60},
            {"kind": "node", "id": "B", "type": "EC2", "x": 300, "y": 230},
            {"kind": "node", "id": "X", "type": "EC2", "x": 300, "y": 420},
        ],
        "links": [link],
    }


_DETOUR = [{"x": 470, "y": 150}, {"x": 470, "y": 360}]


def _root(raw):
    validate(raw)
    diagram = parse_diagram(raw)
    return diagram, build_layout(diagram, REGISTRY)


def test_model_parses_waypoints():
    diagram = parse_diagram(_doc({"waypoints": _DETOUR}))
    assert diagram.links[0].waypoints == [(470, 150), (470, 360)]


def test_plain_link_crosses_but_waypoints_route_clear():
    diagram, root = _root(_doc({}))
    assert any("through element 'B'" in w for w in link_crossing_warnings(root, diagram.links, REGISTRY, 0))

    diagram, root = _root(_doc({"waypoints": _DETOUR}))
    assert link_crossing_warnings(root, diagram.links, REGISTRY, 0) == []


def test_render_plan_threads_the_waypoints_between_endpoints():
    diagram, root = _root(_doc({"waypoints": _DETOUR}))
    by_id = {b.element.id: b for b in root.children}
    _s, _e, style, path = link_render_plan(by_id["A"], by_id["X"], diagram.links[0])
    assert style == "straight"
    assert path[1:-1] == [(470.0, 150.0), (470.0, 360.0)]  # vias, in order, between the two endpoints


def test_endpoint_attaches_on_the_side_facing_the_first_waypoint():
    # First via is down-and-right of A, so A should exit its right edge (idx 3).
    diagram, root = _root(_doc({"waypoints": _DETOUR}))
    by_id = {b.element.id: b for b in root.children}
    s_idx, _e, _style, _path = link_render_plan(by_id["A"], by_id["X"], diagram.links[0])
    assert s_idx == 3


def test_explicit_side_overrides_waypoint_facing():
    diagram, root = _root(_doc({"waypoints": _DETOUR, "fromSide": "bottom"}))
    by_id = {b.element.id: b for b in root.children}
    s_idx, _e, _style, _path = link_render_plan(by_id["A"], by_id["X"], diagram.links[0])
    assert s_idx == 2  # bottom, as the author asked, not the auto right-facing


def test_cross_axis_sides_are_fatal_without_waypoints_but_allowed_with():
    with pytest.raises(DiagramError):
        validate(_doc({"fromSide": "bottom", "toSide": "left"}))
    # waypoints make the routing explicit, so any side combination is fine
    validate(_doc({"fromSide": "bottom", "toSide": "left", "waypoints": _DETOUR}))


def test_schema_rejects_malformed_waypoints():
    with pytest.raises(DiagramError):
        validate(_doc({"waypoints": [{"x": 10}]}))  # missing y
    with pytest.raises(DiagramError):
        validate(_doc({"waypoints": []}))  # minItems: 1


def test_label_anchor_matches_chord_midpoint_for_straight_and_elbow():
    # 2-point straight: arc midpoint == chord midpoint.
    assert link_label_anchor([(0.0, 0.0), (100.0, 0.0)]) == (50.0, 0.0)
    # symmetric elbow Z: arc midpoint lands on the chord midpoint too, so
    # existing (non-waypoint) label placement is unchanged.
    elbow = [(0.0, 0.0), (50.0, 0.0), (50.0, 100.0), (100.0, 100.0)]
    ax, ay = link_label_anchor(elbow)
    assert (round(ax, 3), round(ay, 3)) == (50.0, 50.0)


def test_label_anchor_lands_on_the_polyline_for_a_detour():
    # For the detour, the chord midpoint would sit on B (x=300); the arc
    # midpoint must instead sit on the drawn line, out at the via column.
    diagram, root = _root(_doc({"waypoints": _DETOUR, "label": "detour"}))
    by_id = {b.element.id: b for b in root.children}
    _s, _e, _style, path = link_render_plan(by_id["A"], by_id["X"], diagram.links[0])
    ax, _ay = link_label_anchor(path)
    assert ax > 400  # on the detour column, not back on the chord through B


def test_waypoint_link_renders_as_chained_straight_connectors():
    diagram, root = _root(_doc({"waypoints": _DETOUR}))
    prs = render(diagram, root, REGISTRY)
    connectors = [s for s in prs.slides[0].shapes if isinstance(s, Connector)]
    # path is [A-edge, via1, via2, X-edge] == 3 straight segments for the link.
    assert len(connectors) == 3


def test_drawio_export_carries_the_waypoints():
    diagram, root = _root(_doc({"waypoints": _DETOUR}))
    xml = export_drawio(diagram, root, REGISTRY)
    assert '<Array as="points">' in xml
    assert '<mxPoint x="470.00" y="150.00"/>' in xml
    assert '<mxPoint x="470.00" y="360.00"/>' in xml
