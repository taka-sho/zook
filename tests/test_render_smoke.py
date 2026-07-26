from pathlib import Path

import yaml
from pptx import Presentation

from archdiagram.errors import Warnings
from archdiagram.layout import (
    build_layout,
    icon_resolution_warnings,
    link_aliasing_warnings,
    link_crossing_warnings,
    out_of_canvas_warnings,
    overlap_warnings,
)
from archdiagram.model import parse_diagram
from archdiagram.registry import load_registries
from archdiagram.render import render
from archdiagram.validate import validate

FIXTURE = Path(__file__).parent / "fixtures" / "example.yaml"
CLOUD_ACTORS_FIXTURE = Path(__file__).parent / "fixtures" / "example-cloud-actors.yaml"


def _build(raw):
    validate(raw)
    diagram = parse_diagram(raw)
    registry = load_registries()
    root_box = build_layout(diagram, registry)
    warnings = Warnings()
    margin = diagram.canvas.overlap_margin
    for message in icon_resolution_warnings(root_box, registry):
        warnings.add(message)
    for message in out_of_canvas_warnings(root_box, *diagram.canvas.size):
        warnings.add(message)
    for message in overlap_warnings(root_box, registry, margin):
        warnings.add(message)
    for message in link_crossing_warnings(root_box, diagram.links, registry, margin):
        warnings.add(message)
    for message in link_aliasing_warnings(root_box, diagram.links):
        warnings.add(message)
    presentation = render(diagram, root_box, registry)
    return presentation, warnings


def test_example_yaml_renders_without_warnings():
    raw = yaml.safe_load(FIXTURE.read_text())
    presentation, warnings = _build(raw)
    assert warnings.messages == []
    slide = presentation.slides[0]
    # 2 top-level VPC/S3 elements + 3 links (each a connector + 1 label textbox for the labeled link)
    assert len(slide.shapes) >= 5


def test_example_yaml_output_reopens_cleanly(tmp_path):
    raw = yaml.safe_load(FIXTURE.read_text())
    presentation, _ = _build(raw)
    out = tmp_path / "out.pptx"
    presentation.save(out)

    reopened = Presentation(str(out))
    assert reopened.slide_width == presentation.slide_width
    assert len(reopened.slides) == 1


def test_unresolved_icon_produces_warning_not_error():
    raw = {
        "version": "1.0",
        "canvas": {"aspectRatio": "16:9"},
        "elements": [{"kind": "node", "id": "mystery", "type": "NotARealService"}],
    }
    _, warnings = _build(raw)
    assert any("NotARealService" in m for m in warnings.messages)


def test_cli_exits_nonzero_on_fatal_error(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    bad_yaml = tmp_path / "bad.yaml"
    bad_yaml.write_text(
        """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: node
    id: a
    type: EC2
  - kind: node
    id: a
    type: S3
"""
    )
    runner = CliRunner()
    result = runner.invoke(main, ["build", str(bad_yaml), "-o", str(tmp_path / "out.pptx")])
    assert result.exit_code == 1
    assert "Duplicate" in result.output


def test_cli_succeeds_on_example_yaml(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    out_path = tmp_path / "out.pptx"
    runner = CliRunner()
    result = runner.invoke(main, ["build", str(FIXTURE), "-o", str(out_path)])
    assert result.exit_code == 0
    assert out_path.exists()


def test_example_cloud_actors_renders_without_warnings():
    raw = yaml.safe_load(CLOUD_ACTORS_FIXTURE.read_text())
    presentation, warnings = _build(raw)
    assert warnings.messages == []
    slide = presentation.slides[0]
    assert len(slide.shapes) > 0


def test_cli_succeeds_on_example_cloud_actors(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    out_path = tmp_path / "out.pptx"
    runner = CliRunner()
    result = runner.invoke(main, ["build", str(CLOUD_ACTORS_FIXTURE), "-o", str(out_path)])
    assert result.exit_code == 0
    assert out_path.exists()


def test_cli_reports_overlapping_elements_as_warning_not_error(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    overlapping_yaml = tmp_path / "overlap.yaml"
    overlapping_yaml.write_text(
        """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: node
    id: a
    type: EC2
    x: 100
    y: 100
  - kind: node
    id: b
    type: EC2
    x: 105
    y: 105
"""
    )
    runner = CliRunner()
    result = runner.invoke(main, ["build", str(overlapping_yaml), "-o", str(tmp_path / "out.pptx")])
    assert result.exit_code == 0
    assert "overlaps" in (result.output + result.stderr)


def test_cli_reports_link_crossing_as_warning_not_error(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    crossing_yaml = tmp_path / "crossing.yaml"
    crossing_yaml.write_text(
        """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: node
    id: a
    type: EC2
    x: 0
    y: 100
  - kind: node
    id: b
    type: EC2
    x: 150
    y: 100
  - kind: node
    id: c
    type: EC2
    x: 300
    y: 100
links:
  - from: a
    to: c
"""
    )
    runner = CliRunner()
    result = runner.invoke(main, ["build", str(crossing_yaml), "-o", str(tmp_path / "out.pptx")])
    assert result.exit_code == 0
    assert "passes through element 'b'" in (result.output + result.stderr)


def test_cli_reports_link_aliasing_as_warning_not_error(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    # Reproduces the reported "Z-route false edge aliasing": alb->vpc and
    # vpc->s3 both pick VPC's top-center connection point, so their Z-routes
    # meet nose-to-tail and read as one straight alb->s3 line.
    aliasing_yaml = tmp_path / "aliasing.yaml"
    aliasing_yaml.write_text(
        """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: node
    id: alb
    type: ELB
    x: 297
    y: 74
    width: 100
    height: 100
  - kind: container
    id: vpc
    type: vpc
    x: 405
    y: 229
    width: 300
    height: 300
  - kind: node
    id: s3
    type: S3
    x: 713
    y: 74
    width: 100
    height: 100
links:
  - from: alb
    to: vpc
  - from: vpc
    to: s3
"""
    )
    runner = CliRunner()
    result = runner.invoke(main, ["build", str(aliasing_yaml), "-o", str(tmp_path / "out.pptx")])
    assert result.exit_code == 0
    assert "share a collinear segment" in (result.output + result.stderr)


def test_cli_strict_exits_nonzero_on_warning(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    unknown_yaml = tmp_path / "unknown.yaml"
    unknown_yaml.write_text(
        """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: node
    id: a
    type: NotARealService
"""
    )
    runner = CliRunner()
    lenient = runner.invoke(main, ["build", str(unknown_yaml), "-o", str(tmp_path / "out1.pptx")])
    assert lenient.exit_code == 0

    strict = runner.invoke(main, ["build", str(unknown_yaml), "-o", str(tmp_path / "out2.pptx"), "--strict"])
    assert strict.exit_code == 1


def test_cli_validate_does_not_write_a_file(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    runner = CliRunner()
    with runner.isolated_filesystem(temp_dir=tmp_path):
        result = runner.invoke(main, ["validate", str(FIXTURE)])
        assert result.exit_code == 0
        assert not any(Path(".").iterdir())


def test_cli_validate_catches_unresolved_icon_without_rendering():
    from click.testing import CliRunner

    from archdiagram.cli import main

    runner = CliRunner()
    with runner.isolated_filesystem():
        Path("unknown.yaml").write_text(
            """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: node
    id: a
    type: NotARealService
"""
        )
        result = runner.invoke(main, ["validate", "unknown.yaml", "--format", "json"])
        assert result.exit_code == 0
        assert "NotARealService" in result.output


def test_cli_icons_list_includes_builtin_aws_service():
    from click.testing import CliRunner

    from archdiagram.cli import main

    runner = CliRunner()
    result = runner.invoke(main, ["icons", "list", "--provider", "aws"])
    assert result.exit_code == 0
    assert "EC2" in result.output


def test_cli_preview_writes_a_png(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    out_path = tmp_path / "out.png"
    runner = CliRunner()
    result = runner.invoke(main, ["preview", str(FIXTURE), "-o", str(out_path)])
    assert result.exit_code == 0
    assert out_path.exists()
    assert out_path.stat().st_size > 0


def test_custom_size_and_font_sizes_render_and_preview_without_error(tmp_path):
    from click.testing import CliRunner

    from archdiagram.cli import main

    raw = tmp_path / "sizes.yaml"
    raw.write_text(
        """
version: "1.0"
canvas:
  aspectRatio: "16:9"
elements:
  - kind: container
    id: vpc
    type: vpc
    provider: aws
    label: VPC
    style:
      labelFontSize: 24
    children:
      - kind: node
        id: a
        type: EC2
        size: 100
        style:
          labelFontSize: 18
      - kind: node
        id: b
        type: S3
        size: 40
links:
  - from: a
    to: b
    label: talks to
    labelFontSize: 16
"""
    )
    runner = CliRunner()
    pptx_out = tmp_path / "out.pptx"
    result = runner.invoke(main, ["build", str(raw), "-o", str(pptx_out)])
    assert result.exit_code == 0
    assert pptx_out.exists()

    png_out = tmp_path / "out.png"
    result = runner.invoke(main, ["preview", str(raw), "-o", str(png_out)])
    assert result.exit_code == 0
    assert png_out.exists()
    assert png_out.stat().st_size > 0
