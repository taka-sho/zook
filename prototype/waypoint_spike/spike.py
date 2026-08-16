"""Technical spike: can we render an arbitrary polyline (waypoint) connector
in a .pptx via python-pptx, with an arrowhead, that renders correctly?

The current renderer (src/zook/render.py) draws links with python-pptx's
built-in connector *presets* (STRAIGHT / ELBOW=bentConnector3 / CURVE). Those
presets bend at most once, so they can't route through author/doctor-supplied
waypoints. This spike builds the same little scene three different ways and
writes one .pptx per approach; render each to PNG with LibreOffice to compare.

Scene (logical px): A at top, X at bottom, obstacle B dead between them on the
same vertical line. A single connector must go A -> (detour right of B) -> X
with an arrowhead at X. That's exactly the case doctor stage 3 can't fix when
everything is author-pinned, and the one waypoints would unlock.

    A
    |
    B      link A->X must detour around B
    |
    X

Approaches:
  1. chained  - one STRAIGHT connector per segment; arrow on the last only.
                Reuses the exact primitive render.py already relies on.
  2. freeform - a single freeform (custGeom) shape through all points; arrow
                on its line. One shape, but a new primitive with its own quirks.
  3. custgeom - a real connector (cxnSp) whose preset geometry we replace with
                a hand-written multi-segment custGeom path; arrow on its line.
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU = 9525  # 1 logical px @ 96dpi
LINE = "545B64"
OUT = Path(__file__).parent

# The scene: three boxes and the detour path (logical px).
BOX = 64
A = (300, 60)
B = (300, 230)
X = (300, 420)
# A-bottom -> down -> right of B -> down -> X-top
PATH = [(332, 124), (332, 177), (470, 177), (470, 360), (332, 360), (332, 420)]


def _px(v):
    return Emu(int(v * EMU))


def _new_slide():
    prs = Presentation()
    prs.slide_width = _px(960)
    prs.slide_height = _px(540)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _add_boxes(slide):
    for (cx, cy), label in ((A, "A"), (B, "B (obstacle)"), (X, "X")):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _px(cx), _px(cy), _px(BOX), _px(BOX))
        sp.text_frame.text = label
        sp.text_frame.paragraphs[0].font.size = Pt(9)


def _arrow(line):
    ln = line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def approach_chained():
    from pptx.dml.color import RGBColor

    prs, slide = _new_slide()
    _add_boxes(slide)
    conns = []
    for (x0, y0), (x1, y1) in zip(PATH, PATH[1:]):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _px(x0), _px(y0), _px(x1), _px(y1))
        c.begin_x, c.begin_y = _px(x0), _px(y0)
        c.end_x, c.end_y = _px(x1), _px(y1)
        c.line.color.rgb = RGBColor.from_string(LINE)
        c.line.width = Pt(1.25)
        conns.append(c)
    _arrow(conns[-1].line)  # single arrowhead, on the final segment
    prs.save(OUT / "approach1_chained.pptx")


def approach_freeform():
    from pptx.dml.color import RGBColor

    prs, slide = _new_slide()
    _add_boxes(slide)
    fb = slide.shapes.build_freeform(PATH[0][0], PATH[0][1], scale=EMU)
    fb.add_line_segments([(x, y) for x, y in PATH[1:]], close=False)
    shape = fb.convert_to_shape()
    shape.fill.background()
    shape.line.color.rgb = RGBColor.from_string(LINE)
    shape.line.width = Pt(1.25)
    _arrow(shape.line)
    prs.save(OUT / "approach2_freeform.pptx")


def approach_custgeom():
    from pptx.dml.color import RGBColor

    prs, slide = _new_slide()
    _add_boxes(slide)
    # Start from a real connector so we get a <p:cxnSp>, then swap its preset
    # <a:prstGeom> for a hand-built multi-segment <a:custGeom>.
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _px(PATH[0][0]), _px(PATH[0][1]), _px(PATH[-1][0]), _px(PATH[-1][1]))
    c.line.color.rgb = RGBColor.from_string(LINE)
    c.line.width = Pt(1.25)
    _arrow(c.line)

    xs = [p[0] for p in PATH]
    ys = [p[1] for p in PATH]
    ox, oy = min(xs), min(ys)
    w, h = max(xs) - ox, max(ys) - oy

    spPr = c._element.spPr
    # normalise the shape's own offset/extent to the path bounding box
    xfrm = spPr.find(qn("a:xfrm"))
    xfrm.find(qn("a:off")).set("x", str(int(ox * EMU)))
    xfrm.find(qn("a:off")).set("y", str(int(oy * EMU)))
    xfrm.find(qn("a:ext")).set("cx", str(int(w * EMU)))
    xfrm.find(qn("a:ext")).set("cy", str(int(h * EMU)))

    old = spPr.find(qn("a:prstGeom"))
    spPr.remove(old)
    cust = spPr.makeelement(qn("a:custGeom"), {})
    for tag in ("a:avLst", "a:gdLst", "a:ahLst", "a:cxnLst", "a:rect"):
        cust.append(cust.makeelement(qn(tag), {} if tag != "a:rect" else {"l": "0", "t": "0", "r": str(int(w * EMU)), "b": str(int(h * EMU))}))
    pathLst = cust.makeelement(qn("a:pathLst"), {})
    pth = pathLst.makeelement(qn("a:path"), {"w": str(int(w * EMU)), "h": str(int(h * EMU))})

    def pt(px, py):
        e = pth.makeelement(qn("a:pt"), {"x": str(int((px - ox) * EMU)), "y": str(int((py - oy) * EMU))})
        return e

    moveTo = pth.makeelement(qn("a:moveTo"), {})
    moveTo.append(pt(*PATH[0]))
    pth.append(moveTo)
    for px, py in PATH[1:]:
        lnTo = pth.makeelement(qn("a:lnTo"), {})
        lnTo.append(pt(px, py))
        pth.append(lnTo)
    pathLst.append(pth)
    cust.append(pathLst)
    # custGeom must sit where prstGeom was: after xfrm, before <a:ln>
    ln = spPr.find(qn("a:ln"))
    spPr.insert(list(spPr).index(ln), cust)
    prs.save(OUT / "approach3_custgeom.pptx")


if __name__ == "__main__":
    approach_chained()
    approach_freeform()
    approach_custgeom()
    print("wrote approach1_chained.pptx, approach2_freeform.pptx, approach3_custgeom.pptx")
