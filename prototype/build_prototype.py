"""Prototype to resolve the three open handoff items in detailed-design-pptx.md sec8.6:

1. Connection point index assignment order for rectangular shapes.
2. Actual on-screen behavior of a connector label (txBody injection vs. midpoint textbox).
3. Appropriate rasterization DPI for SVG->PNG icons.

Run:  ../.venv/bin/python build_prototype.py
Output: prototype_output.pptx (then converted to PNG via LibreOffice headless for visual review)
"""

import io
import os

import cairosvg
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))

# Confirmed exact: 1280x720 (16:9) slide is 13.333in x 7.5in, and 914400 EMU/in * 13.333 / 1280
# == 9525 EMU per logical unit. Same factor holds for the 4:3 canvas (960x720 -> 10in x 7.5in).
# I.e. 1 logical unit == 1 px at 96 DPI, exactly, for both supported aspect ratios.
LOGICAL_TO_EMU = 9525


def E(v):
    return Emu(int(round(v * LOGICAL_TO_EMU)))


def add_container_rect(shapes, x, y, w, h, label, border_hex="8C4FFF", dashed=False):
    from pptx.enum.shapes import MSO_SHAPE

    box = shapes.add_shape(MSO_SHAPE.RECTANGLE, E(x), E(y), E(w), E(h))
    box.fill.background()
    box.line.color.rgb = RGBColor.from_string(border_hex)
    box.line.width = Pt(1.25)
    if dashed:
        ln = box.line._get_or_add_ln()
        from pptx.oxml.ns import qn

        prstDash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(prstDash)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor.from_string(border_hex)
    tf.word_wrap = True
    return box


def add_icon_node(shapes, x, y, size, label, png_bytes):
    pic = shapes.add_picture(io.BytesIO(png_bytes), E(x), E(y), E(size), E(size))
    tb = shapes.add_textbox(E(x - 20), E(y + size + 4), E(size + 40), E(20))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(9)
    return pic, tb


def choose_connection_indices(from_box, to_box):
    """Design-doc sec3.3 heuristic: pick the edge that faces the other shape.

    Confirmed via python-pptx source (pptx/shapes/connector.py):
      idx 0 = top-center, 1 = left-center, 2 = bottom-center, 3 = right-center
      (top-center first, proceeding counter-clockwise).
    """
    fx, fy, fw, fh = from_box
    tx, ty, tw, th = to_box
    fcx, fcy = fx + fw / 2, fy + fh / 2
    tcx, tcy = tx + tw / 2, ty + th / 2
    dx, dy = tcx - fcx, tcy - fcy
    if abs(dx) >= abs(dy):
        return (3, 1) if dx >= 0 else (1, 3)
    return (2, 0) if dy >= 0 else (0, 2)


def add_link(shapes, from_shape, from_box, to_shape, to_box, label=None, arrow="end", style="straight"):
    start_idx, end_idx = choose_connection_indices(from_box, to_box)
    connector_type = MSO_CONNECTOR.ELBOW if style == "elbow" else MSO_CONNECTOR.STRAIGHT
    conn = shapes.add_connector(connector_type, E(0), E(0), E(1), E(1))
    conn.begin_connect(from_shape, start_idx)
    conn.end_connect(to_shape, end_idx)
    conn.line.color.rgb = RGBColor.from_string("545B64")
    conn.line.width = Pt(1.25)
    if arrow != "none":
        from pptx.oxml.ns import qn

        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
        if arrow == "both":
            head = ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"})
            ln.append(head)

    if label:
        # sec8.3 finding: p:cxnSp cannot contain a:txBody per the OOXML schema
        # (confirmed: LibreOffice oox source explicitly rejects txBody on cxnSp).
        # So the *only* viable option is a separate textbox at the connector midpoint.
        mx = (conn.begin_x + conn.end_x) / 2
        my = (conn.begin_y + conn.end_y) / 2
        lbl_w, lbl_h = E(60), E(18)
        tb = shapes.add_textbox(Emu(int(mx - lbl_w / 2)), Emu(int(my - lbl_h / 2)), lbl_w, lbl_h)
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        tb.line.fill.background()
        tf = tb.text_frame
        tf.margin_top = tf.margin_bottom = 0
        tf.paragraphs[0].text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(8)
    return conn


def rasterize(size_px):
    svg_path = os.path.join(HERE, "icon_placeholder.svg")
    return cairosvg.svg2png(url=svg_path, output_width=size_px, output_height=size_px)


def build_slide1_hierarchy_and_connectors(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    icon64 = rasterize(64)  # native 1x raster, matches on-slide display size

    # --- VPC group ---
    vpc_group = shapes.add_group_shape()
    vpc_shapes = vpc_group.shapes
    add_container_rect(vpc_shapes, 40, 40, 1120, 420, "Production VPC", "8C4FFF")

    # --- AZ-A group (nested inside VPC group) ---
    az_a_group = vpc_shapes.add_group_shape()
    az_a_shapes = az_a_group.shapes
    add_container_rect(az_a_shapes, 80, 100, 480, 320, "ap-northeast-1a", "00A4A6", dashed=True)
    web_a_box = (120, 170, 64, 64)
    db_a_box = (120, 300, 64, 64)
    web_a, _ = add_icon_node(az_a_shapes, *web_a_box[:2], 64, "WebServer A", icon64)
    db_a, _ = add_icon_node(az_a_shapes, *db_a_box[:2], 64, "Primary DB", icon64)
    add_link(az_a_shapes, web_a, web_a_box, db_a, db_a_box, label="3306")

    # --- AZ-C group (nested inside VPC group) ---
    az_c_group = vpc_shapes.add_group_shape()
    az_c_shapes = az_c_group.shapes
    add_container_rect(az_c_shapes, 600, 100, 480, 320, "ap-northeast-1c", "00A4A6", dashed=True)
    web_c_box = (640, 170, 64, 64)
    fn_c_box = (900, 170, 64, 64)
    web_c, _ = add_icon_node(az_c_shapes, *web_c_box[:2], 64, "WebServer C", icon64)
    fn_c, _ = add_icon_node(az_c_shapes, *fn_c_box[:2], 64, "Batch Worker", icon64)
    add_link(az_c_shapes, web_c, web_c_box, fn_c, fn_c_box, style="elbow")

    # Node outside the VPC, absolute position (like example.yaml's S3 bucket)
    bucket_box = (1150, 300, 64, 64)
    bucket, _ = add_icon_node(shapes, *bucket_box[:2], 64, "Asset Bucket", icon64)
    add_link(shapes, web_c, web_c_box, bucket, bucket_box, arrow="none")

    # --- Connection-index legend: visually confirms idx 0/1/2/3 = top/left/bottom/right ---
    legend_box = (150, 500, 100, 100)
    legend = add_container_rect(shapes, *legend_box, "idx legend", "545B64")
    targets = {
        0: (150 + 50 - 8, 460, 16, 16),   # above -> top
        1: (60, 500 + 50 - 8, 16, 16),    # left -> left
        2: (150 + 50 - 8, 620, 16, 16),   # below -> bottom
        3: (270, 500 + 50 - 8, 16, 16),   # right -> right
    }
    for idx, (tx, ty, tw, th) in targets.items():
        marker = shapes.add_shape(9, E(tx), E(ty), E(tw), E(th))  # 9 = OVAL
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor.from_string("D13212")
        marker.line.fill.background()
        marker.text_frame.text = str(idx)
        marker.text_frame.paragraphs[0].font.size = Pt(8)
        marker.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("FFFFFF")
        marker.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, E(0), E(0), E(1), E(1))
        conn.begin_connect(legend, idx)
        mcx, mcy = tx + tw / 2, ty + th / 2
        conn.end_x, conn.end_y = E(mcx), E(mcy)
        conn.line.color.rgb = RGBColor.from_string("D13212")

    cap = shapes.add_textbox(E(40), E(650), E(1200), E(60))
    cap.text_frame.text = (
        "Legend: red dot N marks connection-point index N on the gray box "
        "(0=top, 1=left, 2=bottom, 3=right per pptx/shapes/connector.py). "
        "Link labels are separate textboxes at the connector midpoint "
        "(cxnSp cannot carry txBody per OOXML schema)."
    )
    cap.text_frame.paragraphs[0].font.size = Pt(10)
    cap.text_frame.word_wrap = True


def build_slide2_dpi_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    shapes.add_textbox(E(40), E(30), E(1200), E(30)).text_frame.text = (
        "Icon DPI comparison — same on-slide display size (80x80 logical units), "
        "different source PNG pixel resolution"
    )

    display_size = 80  # logical units, matches a slightly-larger-than-default node
    cases = [
        (96, "96 dpi-equiv (1x, 80px source)"),
        (192, "192 dpi-equiv (2x, 160px source)"),
        (288, "288 dpi-equiv (3x, 240px source)"),
        (384, "384 dpi-equiv (4x, 320px source)"),
    ]
    x = 100
    for dpi, caption in cases:
        px = int(round(display_size * LOGICAL_TO_EMU / 914400 * dpi))
        png = rasterize(px)
        shapes.add_picture(io.BytesIO(png), E(x), E(150), E(display_size), E(display_size))
        cap = shapes.add_textbox(E(x - 40), E(240), E(display_size + 80), E(60))
        cap.text_frame.text = f"{caption}\n{px}x{px}px source"
        cap.text_frame.word_wrap = True
        for p in cap.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9)
        x += 220


def main():
    prs = Presentation()
    prs.slide_width = E(1280)
    prs.slide_height = E(720)
    build_slide1_hierarchy_and_connectors(prs)
    build_slide2_dpi_comparison(prs)
    out_path = os.path.join(HERE, "prototype_output.pptx")
    prs.save(out_path)
    print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
